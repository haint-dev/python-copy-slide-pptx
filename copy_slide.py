#!/usr/bin/env python3
"""
Copy slides from multiple source PPTX files into a template-based presentation.

Usage:
    Flexible API to:
    1. Create a new presentation from a template (keeping theme/layouts, removing slides)
    2. Copy selected slides from any number of source files
    3. Apply the template's layout/background to copied slides
"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from lxml import etree
from copy import deepcopy
import os
import re

# Counter for generating unique media filenames
_media_counter = 0


def create_from_template(template_path):
    """
    Create a new empty presentation that inherits all themes/layouts from the template.
    All existing slides are removed.

    Args:
        template_path: Path to the template PPTX file

    Returns:
        A Presentation object with no slides but all template themes/layouts intact
    """
    prs = Presentation(template_path)

    # Remove all existing slides (keep themes/layouts/masters)
    sldIdLst = prs.element.find(qn('p:sldIdLst'))
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    return prs


def _next_media_partname(content_type):
    """Generate a unique partname for a media file."""
    global _media_counter
    _media_counter += 1
    ext_map = {
        'image/png': '.png',
        'image/jpeg': '.jpg',
        'image/gif': '.gif',
        'image/bmp': '.bmp',
        'image/tiff': '.tiff',
        'image/svg+xml': '.svg',
        'image/x-emf': '.emf',
        'image/x-wmf': '.wmf',
    }
    ext = ext_map.get(content_type, '.bin')
    return PackURI(f'/ppt/media/copied_image{_media_counter}{ext}')


def _copy_relationships(src_slide, dst_slide):
    """
    Copy all non-layout, non-notes relationships from source slide to destination.
    Creates new Part copies for media to avoid duplicate partname conflicts.
    Returns a mapping of old rId -> new rId.
    """
    rid_map = {}

    for rel in src_slide.part.rels.values():
        # Skip layout and notes relationships (destination has its own)
        if 'slideLayout' in rel.reltype or 'notesSlide' in rel.reltype:
            continue

        if rel.is_external:
            # External relationships (hyperlinks, etc.)
            new_rid = dst_slide.part.rels.get_or_add_ext_rel(
                rel.reltype, rel.target_ref
            )
            rid_map[rel.rId] = new_rid
        else:
            # Internal relationships (images, charts, embedded objects, etc.)
            # Create a fresh Part copy with unique name to avoid ZIP duplicate warnings
            src_part = rel.target_part
            if 'image' in rel.reltype:
                new_partname = _next_media_partname(src_part.content_type)
                new_part = Part(
                    new_partname, src_part.content_type,
                    dst_slide.part.package, src_part.blob
                )
                new_rid = dst_slide.part.relate_to(new_part, rel.reltype)
            else:
                new_rid = dst_slide.part.relate_to(src_part, rel.reltype)
            rid_map[rel.rId] = new_rid

    return rid_map


def _update_rids_in_tree(tree, rid_map):
    """Replace old rId references with new ones throughout an XML tree.
    Uses single-pass replacement to avoid cascading (e.g. rId4->rId5 then rId5->rId4).
    """
    for elem in tree.iter():
        for attr_name in list(elem.attrib.keys()):
            val = elem.get(attr_name)
            if val in rid_map and rid_map[val] != val:
                elem.set(attr_name, rid_map[val])


def _is_title_shape(sp_elem):
    """Check if a shape element is a title placeholder."""
    info = _get_placeholder_info(sp_elem)
    return info is not None and info[0] in ('title', 'ctrTitle')


def _get_placeholder_info(sp_elem):
    """
    Get placeholder info from a shape element.
    Returns (type, idx) tuple if shape is a placeholder, None otherwise.
    """
    if sp_elem.tag != qn('p:sp'):
        return None
    nvSpPr = sp_elem.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return None
    nvPr = nvSpPr.find(qn('p:nvPr'))
    if nvPr is None:
        return None
    ph = nvPr.find(qn('p:ph'))
    if ph is None:
        return None
    return (ph.get('type', ''), ph.get('idx', ''))


def _copy_placeholder_content(src_sp, dst_sp):
    """
    Copy text paragraphs from source placeholder to destination placeholder.
    Keeps destination's bodyPr and lstStyle (from template layout).
    """
    src_txBody = src_sp.find(qn('p:txBody'))
    dst_txBody = dst_sp.find(qn('p:txBody'))

    if src_txBody is None or dst_txBody is None:
        return

    # Remove existing paragraphs from destination
    for p in list(dst_txBody.findall(qn('a:p'))):
        dst_txBody.remove(p)

    # Copy paragraphs from source
    for p in src_txBody.findall(qn('a:p')):
        dst_txBody.append(deepcopy(p))


def _remove_placeholder_ref(sp_elem):
    """Remove placeholder reference from a shape so it becomes a regular shape."""
    nvSpPr = sp_elem.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return
    nvPr = nvSpPr.find(qn('p:nvPr'))
    if nvPr is None:
        return
    ph = nvPr.find(qn('p:ph'))
    if ph is not None:
        nvPr.remove(ph)


def _map_placeholders(src_slide, dst_slide, rid_map):
    """
    Map content from source placeholders into template placeholders.
    Non-placeholder shapes and unmatched placeholders are added as extra shapes.
    """
    shape_tags = {qn('p:sp'), qn('p:grpSp'), qn('p:pic'),
                  qn('p:graphicFrame'), qn('p:cxnSp')}

    # Classify source shapes
    src_ph_map = {}     # type_key -> sp_elem
    src_other = []      # non-placeholder shapes

    for child in src_slide.shapes._spTree:
        if child.tag not in shape_tags:
            continue
        ph_info = _get_placeholder_info(child)
        if ph_info:
            type_key = ph_info[0] or f'_idx_{ph_info[1]}'
            if type_key not in src_ph_map:
                src_ph_map[type_key] = child
        else:
            src_other.append(child)

    # Map content to matching template placeholders
    dst_spTree = dst_slide.shapes._spTree
    matched = set()

    for dst_child in list(dst_spTree):
        dst_ph_info = _get_placeholder_info(dst_child)
        if dst_ph_info is None:
            continue
        dst_key = dst_ph_info[0] or f'_idx_{dst_ph_info[1]}'
        if dst_key in src_ph_map and dst_key not in matched:
            _copy_placeholder_content(src_ph_map[dst_key], dst_child)
            matched.add(dst_key)

    # Add non-placeholder shapes from source
    for sp in src_other:
        new_sp = deepcopy(sp)
        _update_rids_in_tree(new_sp, rid_map)
        dst_spTree.append(new_sp)

    # Add unmatched source placeholders as regular shapes
    for type_key, src_sp in src_ph_map.items():
        if type_key not in matched:
            new_sp = deepcopy(src_sp)
            _update_rids_in_tree(new_sp, rid_map)
            _remove_placeholder_ref(new_sp)
            dst_spTree.append(new_sp)


def _extract_theme_colors(prs):
    """
    Extract the color scheme from a presentation's theme.
    Returns dict mapping scheme name -> hex value (uppercase, e.g. 'accent1' -> 'FF5030').
    """
    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_elem = etree.fromstring(rel.target_part.blob)
                clrScheme = theme_elem.find('.//' + qn('a:clrScheme'))
                if clrScheme is None:
                    continue
                colors = {}
                for name in ['dk1', 'dk2', 'lt1', 'lt2',
                             'accent1', 'accent2', 'accent3', 'accent4',
                             'accent5', 'accent6', 'hlink', 'folHlink']:
                    elem = clrScheme.find(qn(f'a:{name}'))
                    if elem is None:
                        continue
                    srgb = elem.find(qn('a:srgbClr'))
                    if srgb is not None:
                        colors[name] = srgb.get('val', '').upper()
                    else:
                        sys_clr = elem.find(qn('a:sysClr'))
                        if sys_clr is not None:
                            colors[name] = sys_clr.get('lastClr', '').upper()
                if colors:
                    return colors
    return {}


def _remap_colors_to_theme(spTree, src_theme_colors):
    """
    Replace hardcoded srgbClr values that match a source theme color
    with the corresponding schemeClr reference.

    This way, colors that were theme-derived in the source will adapt
    to the destination template's theme automatically.
    Only exact matches are replaced — other hardcoded colors are left as-is.
    """
    # Build reverse mapping: hex value -> first matching scheme name
    hex_to_scheme = {}
    for name, hex_val in src_theme_colors.items():
        if hex_val and hex_val not in hex_to_scheme:
            hex_to_scheme[hex_val] = name

    if not hex_to_scheme:
        return

    for srgb_elem in list(spTree.iter(qn('a:srgbClr'))):
        hex_val = srgb_elem.get('val', '').upper()
        if hex_val not in hex_to_scheme:
            continue

        scheme_name = hex_to_scheme[hex_val]
        parent = srgb_elem.getparent()

        # Create schemeClr element with the same namespace
        scheme_elem = parent.makeelement(qn('a:schemeClr'), {'val': scheme_name})

        # Preserve child modifiers (alpha, tint, shade, lumMod, lumOff, etc.)
        for child in list(srgb_elem):
            scheme_elem.append(child)

        parent.replace(srgb_elem, scheme_elem)


def _remap_fonts_to_theme(spTree):
    """
    Replace all hardcoded font references with theme font references.

    Title shapes -> major font (+mj-lt/ea/cs)
    All other shapes/tables -> minor font (+mn-lt/ea/cs)

    This ensures the output uses the template's theme fonts and avoids
    missing font issues on machines that don't have the source fonts.
    """
    font_tags = {
        qn('a:latin'): '-lt',
        qn('a:ea'): '-ea',
        qn('a:cs'): '-cs',
    }
    metric_attrs = ['panose', 'pitchFamily', 'charset']

    def _remap_element(font_elem, prefix):
        typeface = font_elem.get('typeface', '')
        if typeface.startswith('+'):
            return  # Already a theme reference
        suffix = font_tags.get(font_elem.tag)
        if suffix is None:
            return
        font_elem.set('typeface', prefix + suffix)
        for attr in metric_attrs:
            if attr in font_elem.attrib:
                del font_elem.attrib[attr]

    # Track processed elements to avoid double-processing
    processed = set()

    # Process p:sp shapes with title detection
    for sp in spTree.iter(qn('p:sp')):
        prefix = '+mj' if _is_title_shape(sp) else '+mn'
        for tag in font_tags:
            for font_elem in sp.iter(tag):
                _remap_element(font_elem, prefix)
                processed.add(id(font_elem))

    # Process remaining font elements (tables in graphicFrame, etc.) as minor
    for tag in font_tags:
        for font_elem in spTree.iter(tag):
            if id(font_elem) not in processed:
                _remap_element(font_elem, '+mn')


def copy_slide(src_prs, src_slide_index, dst_prs, layout_index=0, apply_template_bg=True,
               remap_fonts=True, remap_colors=True, src_theme_colors=None,
               use_placeholders=False):
    """
    Copy a single slide from source presentation to destination.

    Args:
        src_prs: Source Presentation object
        src_slide_index: 0-based index of the slide to copy
        dst_prs: Destination Presentation object
        layout_index: Index of the slide layout in destination to apply (0-based)
        apply_template_bg: If True, use template background; if False, copy source background
        use_placeholders: If True, map source content into template placeholders
            instead of replacing the entire shape tree

    Returns:
        The newly created slide in the destination presentation
    """
    src_slide = list(src_prs.slides)[src_slide_index]

    # Create new slide with template layout
    dst_layout = dst_prs.slide_layouts[layout_index]
    dst_slide = dst_prs.slides.add_slide(dst_layout)

    # Step 1: Copy relationships (images, charts, etc.) and build rId mapping
    rid_map = _copy_relationships(src_slide, dst_slide)

    # Step 2: Transfer shapes
    if use_placeholders:
        # Map content into template placeholders; add other shapes alongside
        _map_placeholders(src_slide, dst_slide, rid_map)
    else:
        # Replace entire shape tree with source shapes
        new_spTree = deepcopy(src_slide.shapes._spTree)
        _update_rids_in_tree(new_spTree, rid_map)
        dst_spTree = dst_slide.shapes._spTree
        dst_spTree.getparent().replace(dst_spTree, new_spTree)

    # Step 3: Remap fonts to template theme fonts
    if remap_fonts:
        _remap_fonts_to_theme(dst_slide.shapes._spTree)

    # Step 4: Remap hardcoded colors to theme color references
    if remap_colors and src_theme_colors:
        _remap_colors_to_theme(dst_slide.shapes._spTree, src_theme_colors)

    # Step 5: Handle background
    if not apply_template_bg:
        # Copy source slide's background (if it has one)
        src_cSld = src_slide.element.find(qn('p:cSld'))
        src_bg = src_cSld.find(qn('p:bg')) if src_cSld is not None else None
        if src_bg is not None:
            dst_cSld = dst_slide.element.find(qn('p:cSld'))
            new_bg = deepcopy(src_bg)
            # Insert background before spTree
            dst_spTree_elem = dst_cSld.find(qn('p:spTree'))
            dst_cSld.insert(list(dst_cSld).index(dst_spTree_elem), new_bg)

    return dst_slide


def copy_slides_to_template(template_path, slide_selections, output_path,
                            layout_index=0, apply_template_bg=True,
                            remap_fonts=True, remap_colors=True,
                            use_placeholders=False):
    """
    Main function: create presentation from template and copy selected slides.

    Args:
        template_path: Path to the template PPTX file
        slide_selections: List of tuples (source_path, slide_indices)
            - source_path: Path to a source PPTX file
            - slide_indices: List of 0-based slide indices to copy
        output_path: Path for the output PPTX file
        layout_index: Template layout index to apply (0-based)
        apply_template_bg: If True, template background is applied to all slides

    Returns:
        Path to the saved output file
    """
    print(f"Creating presentation from template: {os.path.basename(template_path)}")
    dst_prs = create_from_template(template_path)
    print(f"  Available layouts: {[l.name for l in dst_prs.slide_layouts]}")
    print(f"  Using layout [{layout_index}]: \"{dst_prs.slide_layouts[layout_index].name}\"")
    print()

    slide_count = 0
    for src_path, indices in slide_selections:
        src_prs = Presentation(src_path)
        total = len(list(src_prs.slides))
        src_theme_colors = _extract_theme_colors(src_prs) if remap_colors else None
        print(f"Source: {os.path.basename(src_path)} ({total} slides)")

        for idx in indices:
            if idx < 0 or idx >= total:
                print(f"  [SKIP] Slide {idx} out of range (0-{total-1})")
                continue

            copy_slide(src_prs, idx, dst_prs, layout_index, apply_template_bg,
                       remap_fonts, remap_colors, src_theme_colors, use_placeholders)
            slide_count += 1
            print(f"  Copied slide {idx} -> destination slide {slide_count}")

    dst_prs.save(output_path)
    print(f"\nSaved {slide_count} slides to: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# Helper functions for building slide index lists
# ---------------------------------------------------------------------------

def first_n(pptx_path, n):
    """Return indices of the first N slides."""
    total = len(list(Presentation(pptx_path).slides))
    return list(range(min(n, total)))


def last_n(pptx_path, n):
    """Return indices of the last N slides."""
    total = len(list(Presentation(pptx_path).slides))
    return list(range(max(0, total - n), total))


def slide_range(start, end):
    """Return indices from start to end (inclusive, 0-based)."""
    return list(range(start, end + 1))


# ---------------------------------------------------------------------------
# Example / Demo
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

    template_path = os.path.join(BASE_DIR, 'template.pptx')
    upload_path = os.path.join(BASE_DIR, 'file_upload.pptx')
    output_path = os.path.join(BASE_DIR, 'output.pptx')

    # Task: copy first 5 slides + slide 6 + last 2 slides from file_upload
    slide_selections = [
        (upload_path, first_n(upload_path, 6)),   # Slides 0,1,2,3,4,thực tế đã apply được những gì từ file template nhỉ5
        #(upload_path, [5]),                        # Slide 6 (index 5, 0-based)
        #(upload_path, last_n(upload_path, 2)),     # Slides 9,10
    ]

    copy_slides_to_template(
        template_path=template_path,
        slide_selections=slide_selections,
        output_path=output_path,
        layout_index=0,           # "Cover slide layout"
        apply_template_bg=True,   # Apply template background
    )
