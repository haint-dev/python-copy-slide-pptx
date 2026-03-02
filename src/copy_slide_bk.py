#!/usr/bin/env python3
"""
Copy slides from multiple source PPTX files into a template-based presentation.

Usage:
    Flexible API to:
    1. Create a new presentation from a template (keeping theme/layouts, removing slides)
    2. Copy selected slides from any number of source files
    3. Apply the template's layout/background to copied slides
"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from lxml import etree
from copy import deepcopy
import os
import re

# Counter for generating unique media filenames
_media_counter = 0


_A4_WIDTH = 10692000   # 297mm in EMU (landscape)
_A4_HEIGHT = 7560000   # 210mm in EMU (landscape)


def _scale_template_shapes(prs, old_width, old_height):
    """
    Scale all shapes in slide layouts and slide masters to match the new
    slide dimensions.  Called when slide_size is changed so that template
    decoration shapes (background images, bars, etc.) stretch to fill the
    new canvas instead of leaving gaps.
    """
    new_width = prs.slide_width
    new_height = prs.slide_height
    if new_width == old_width and new_height == old_height:
        return

    sx = new_width / old_width
    sy = new_height / old_height

    for layout in prs.slide_layouts:
        cSld = layout.element.find(qn('p:cSld'))
        if cSld is not None:
            spTree = cSld.find(qn('p:spTree'))
            if spTree is not None:
                _scale_xfrm(spTree, sx, sy)

    for master in prs.slide_masters:
        cSld = master.element.find(qn('p:cSld'))
        if cSld is not None:
            spTree = cSld.find(qn('p:spTree'))
            if spTree is not None:
                _scale_xfrm(spTree, sx, sy)


def create_from_template(template_path, slide_size=None):
    """
    Create a new empty presentation that inherits all themes/layouts from the template.
    All existing slides are removed.

    Args:
        template_path: Path to the template PPTX file
        slide_size: Optional (width, height) in EMU. Use 'a4' for A4 landscape.

    Returns:
        A Presentation object with no slides but all template themes/layouts intact
    """
    prs = Presentation(template_path)

    # Remove all existing slides (keep themes/layouts/masters)
    sldIdLst = prs.element.find(qn('p:sldIdLst'))
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    # Override slide dimensions if requested
    if slide_size is not None:
        old_width = prs.slide_width
        old_height = prs.slide_height
        if slide_size == 'a4':
            prs.slide_width = _A4_WIDTH
            prs.slide_height = _A4_HEIGHT
        else:
            prs.slide_width, prs.slide_height = slide_size
        _scale_template_shapes(prs, old_width, old_height)

    return prs


def _next_media_partname(content_type):
    """Generate a unique partname for a media file."""
    global _media_counter
    _media_counter += 1
    ext_map = {
        'image/png': '.png',
        'image/jpeg': '.jpg',
        'image/gif': '.gif',
        'image/bmp': '.bmp',
        'image/tiff': '.tiff',
        'image/svg+xml': '.svg',
        'image/x-emf': '.emf',
        'image/x-wmf': '.wmf',
    }
    ext = ext_map.get(content_type, '.bin')
    return PackURI(f'/ppt/media/copied_image{_media_counter}{ext}')


def _copy_relationships(src_slide, dst_slide):
    """
    Copy all non-layout, non-notes relationships from source slide to destination.
    Creates new Part copies for media to avoid duplicate partname conflicts.
    Returns a mapping of old rId -> new rId.
    """
    rid_map = {}

    for rel in src_slide.part.rels.values():
        # Skip layout and notes relationships (destination has its own)
        if 'slideLayout' in rel.reltype or 'notesSlide' in rel.reltype:
            continue

        if rel.is_external:
            # External relationships (hyperlinks, etc.)
            new_rid = dst_slide.part.rels.get_or_add_ext_rel(
                rel.reltype, rel.target_ref
            )
            rid_map[rel.rId] = new_rid
        else:
            # Internal relationships (images, charts, embedded objects, etc.)
            # Create a fresh Part copy with unique name to avoid ZIP duplicate warnings
            src_part = rel.target_part
            if 'image' in rel.reltype:
                new_partname = _next_media_partname(src_part.content_type)
                new_part = Part(
                    new_partname, src_part.content_type,
                    dst_slide.part.package, src_part.blob
                )
                new_rid = dst_slide.part.relate_to(new_part, rel.reltype)
            else:
                new_rid = dst_slide.part.relate_to(src_part, rel.reltype)
            rid_map[rel.rId] = new_rid

    return rid_map


def _update_rids_in_tree(tree, rid_map):
    """Replace old rId references with new ones throughout an XML tree.
    Uses single-pass replacement to avoid cascading (e.g. rId4->rId5 then rId5->rId4).
    """
    for elem in tree.iter():
        for attr_name in list(elem.attrib.keys()):
            val = elem.get(attr_name)
            if val in rid_map and rid_map[val] != val:
                elem.set(attr_name, rid_map[val])


def _is_title_shape(sp_elem):
    """Check if a shape element is a title placeholder."""
    info = _get_placeholder_info(sp_elem)
    return info is not None and info[0] in ('title', 'ctrTitle')


def _get_placeholder_info(sp_elem):
    """
    Get placeholder info from a shape element.
    Returns (type, idx) tuple if shape is a placeholder, None otherwise.
    """
    if sp_elem.tag != qn('p:sp'):
        return None
    nvSpPr = sp_elem.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return None
    nvPr = nvSpPr.find(qn('p:nvPr'))
    if nvPr is None:
        return None
    ph = nvPr.find(qn('p:ph'))
    if ph is None:
        return None
    return (ph.get('type', ''), ph.get('idx', ''))


def _copy_placeholder_content(src_sp, dst_sp):
    """
    Copy text paragraphs from source placeholder to destination placeholder.
    Keeps destination's bodyPr and lstStyle (from template layout).
    """
    src_txBody = src_sp.find(qn('p:txBody'))
    dst_txBody = dst_sp.find(qn('p:txBody'))

    if src_txBody is None or dst_txBody is None:
        return

    # Remove existing paragraphs from destination
    for p in list(dst_txBody.findall(qn('a:p'))):
        dst_txBody.remove(p)

    # Copy paragraphs from source
    for p in src_txBody.findall(qn('a:p')):
        dst_txBody.append(deepcopy(p))


def _remove_placeholder_ref(sp_elem):
    """Remove placeholder reference from a shape so it becomes a regular shape."""
    nvSpPr = sp_elem.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return
    nvPr = nvSpPr.find(qn('p:nvPr'))
    if nvPr is None:
        return
    ph = nvPr.find(qn('p:ph'))
    if ph is not None:
        nvPr.remove(ph)


def _find_layout_placeholder(src_slide, ph_type, ph_idx):
    """Find matching placeholder element in source layout or master."""
    src_layout = src_slide.slide_layout
    # Search layout first
    layout_spTree = src_layout.element.find('.//' + qn('p:cSld')).find(qn('p:spTree'))
    for sp in layout_spTree.findall(qn('p:sp')):
        ph = _get_placeholder_info(sp)
        if ph and ph[0] == ph_type and ph[1] == ph_idx:
            return sp
    # Fallback to master
    master_spTree = src_layout.slide_master.element.find('.//' + qn('p:cSld')).find(qn('p:spTree'))
    for sp in master_spTree.findall(qn('p:sp')):
        ph = _get_placeholder_info(sp)
        if ph and ph[0] == ph_type and ph[1] == ph_idx:
            return sp
    return None


def _collect_defRPr_from_layout(layout_sp):
    """
    Collect default run properties (fill, font size) from a layout placeholder.
    Searches lstStyle levels and direct paragraph defRPr.
    Returns (default_fill_elem_or_None, default_sz_str_or_None).
    """
    if layout_sp is None:
        return None, None

    layout_txBody = layout_sp.find(qn('p:txBody'))
    if layout_txBody is None:
        return None, None

    default_fill = None
    default_sz = None

    # Check lstStyle levels
    lstStyle = layout_txBody.find(qn('a:lstStyle'))
    if lstStyle is not None:
        for lvl_tag in ['a:lvl1pPr', 'a:lvl2pPr', 'a:lvl3pPr']:
            lvl = lstStyle.find(qn(lvl_tag))
            if lvl is not None:
                defRPr = lvl.find(qn('a:defRPr'))
                if defRPr is not None:
                    if default_fill is None:
                        fill = defRPr.find(qn('a:solidFill'))
                        if fill is not None:
                            default_fill = deepcopy(fill)
                    if default_sz is None:
                        sz = defRPr.get('sz')
                        if sz:
                            default_sz = sz
                if default_fill is not None and default_sz:
                    break

    # Also check direct defRPr in layout paragraphs
    for p in layout_txBody.findall(qn('a:p')):
        pPr = p.find(qn('a:pPr'))
        if pPr is not None:
            defRPr = pPr.find(qn('a:defRPr'))
            if defRPr is not None:
                if default_fill is None:
                    fill = defRPr.find(qn('a:solidFill'))
                    if fill is not None:
                        default_fill = deepcopy(fill)
                if default_sz is None:
                    sz = defRPr.get('sz')
                    if sz:
                        default_sz = sz
        if default_fill is not None and default_sz:
            break

    return default_fill, default_sz


def _bake_placeholder_styles(sp_elem, src_slide):
    """
    Bake inherited properties from source layout into the shape before
    converting it to a regular shape. Placeholders inherit position, size,
    and text styles from their layout/master. Regular shapes don't inherit,
    so we must resolve and inline these properties.

    Bakes: position/size (xfrm), text color (solidFill), font size (sz).
    For Keynote compatibility, ensures EVERY text run has explicit color
    and font size, falling back to schemeClr tx1 / 1800 (18pt) if no
    layout default is found.
    """
    ph_info = _get_placeholder_info(sp_elem)
    if ph_info is None:
        return

    ph_type, ph_idx = ph_info
    layout_sp = _find_layout_placeholder(src_slide, ph_type, ph_idx)

    # --- Bake position/size ---
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp_elem, qn('p:spPr'))

    if spPr.find(qn('a:xfrm')) is None and layout_sp is not None:
        layout_spPr = layout_sp.find(qn('p:spPr'))
        if layout_spPr is not None:
            xfrm = layout_spPr.find(qn('a:xfrm'))
            if xfrm is not None:
                spPr.insert(0, deepcopy(xfrm))

    # --- Bake default text styles (color, font size) from layout into runs ---
    default_fill, default_sz = _collect_defRPr_from_layout(layout_sp)

    # Fallback: if no explicit fill found from layout, use schemeClr tx1
    # (standard dark text color). This is critical for Keynote which
    # cannot resolve inherited text color for non-placeholder shapes.
    if default_fill is None:
        default_fill = etree.Element(qn('a:solidFill'))
        etree.SubElement(default_fill, qn('a:schemeClr'), {'val': 'tx1'})

    # Fallback font size: 1800 (18pt) is a common default
    if default_sz is None:
        default_sz = '1800'

    # Apply defaults to ALL runs that lack explicit values
    txBody = sp_elem.find(qn('p:txBody'))
    if txBody is None:
        return

    for r in txBody.iter(qn('a:r')):
        rPr = r.find(qn('a:rPr'))
        if rPr is None:
            rPr = r.makeelement(qn('a:rPr'), {})
            r.insert(0, rPr)
        if rPr.find(qn('a:solidFill')) is None and rPr.find(qn('a:noFill')) is None:
            rPr.append(deepcopy(default_fill))
        if not rPr.get('sz'):
            rPr.set('sz', default_sz)


def _get_shape_rect(sp_elem):
    """Extract bounding rectangle (x, y, cx, cy) from a shape's spPr xfrm."""
    spPr = sp_elem.find(qn('p:spPr'))
    if spPr is None:
        return None
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return None
    off = xfrm.find(qn('a:off'))
    ext = xfrm.find(qn('a:ext'))
    if off is None or ext is None:
        return None
    return (int(off.get('x', 0)), int(off.get('y', 0)),
            int(ext.get('cx', 0)), int(ext.get('cy', 0)))


def _rects_overlap(r1, r2):
    """Check if two rectangles overlap. Each rect is (x, y, cx, cy)."""
    return (r1[0] < r2[0] + r2[2] and r1[0] + r1[2] > r2[0] and
            r1[1] < r2[1] + r2[3] and r1[1] + r1[3] > r2[1])


def _find_placeholders_needing_backing(src_slide):
    """
    Find placeholders in the source slide that have text content AND
    inherit a light-colored text fill (bg1, bg2, lt1, lt2) from the layout.

    These placeholders visually depend on a dark decoration shape behind them
    in the source layout.  Returns a list of bounding rectangles (from the
    layout placeholder) that need backing decoration shapes.
    """
    _LIGHT_SCHEMES = {'bg1', 'bg2', 'lt1', 'lt2'}
    regions = []

    for sp in src_slide.shapes._spTree.findall(qn('p:sp')):
        ph_info = _get_placeholder_info(sp)
        if ph_info is None:
            continue

        # Must have text content
        txBody = sp.find(qn('p:txBody'))
        if txBody is None:
            continue
        has_text = any(
            t.text and t.text.strip()
            for t in txBody.iter(qn('a:t'))
        )
        if not has_text:
            continue

        # Check if layout gives this placeholder a light text color
        ph_type, ph_idx = ph_info
        layout_sp = _find_layout_placeholder(src_slide, ph_type, ph_idx)
        if layout_sp is None:
            continue

        default_fill, _ = _collect_defRPr_from_layout(layout_sp)
        if default_fill is None:
            continue

        scheme_clr = default_fill.find(qn('a:schemeClr'))
        if scheme_clr is None or scheme_clr.get('val') not in _LIGHT_SCHEMES:
            continue

        # Get bounding rect from layout placeholder (slide placeholder often
        # has empty spPr and inherits position from layout)
        rect = _get_shape_rect(layout_sp)
        if rect is not None:
            regions.append(rect)

    return regions


def _scale_xfrm(shape_elem, sx, sy):
    """
    Scale position and size of a shape's xfrm by (sx, sy) ratios.
    Handles both direct spPr/xfrm and nested xfrm in group shapes.
    """
    for xfrm in shape_elem.iter(qn('a:xfrm')):
        off = xfrm.find(qn('a:off'))
        ext = xfrm.find(qn('a:ext'))
        if off is not None:
            off.set('x', str(int(int(off.get('x', 0)) * sx)))
            off.set('y', str(int(int(off.get('y', 0)) * sy)))
        if ext is not None:
            ext.set('cx', str(int(int(ext.get('cx', 0)) * sx)))
            ext.set('cy', str(int(int(ext.get('cy', 0)) * sy)))


def _copy_layout_decorations(src_slide, dst_slide, src_prs, dst_prs):
    """
    Copy layout decoration shapes that serve as backgrounds for placeholders
    with light-colored text.

    Only copies decoration shapes that spatially overlap with placeholders
    whose inherited text color is a light scheme color (bg1, bg2, lt1, lt2).
    This ensures dark backing bars are preserved without copying unrelated
    layout decorations (logos, gradient lines, etc.) that would clash with
    the destination template.

    Shapes are scaled proportionally when source and destination slide
    dimensions differ (e.g. 4:3 source → 16:9 template).
    """
    backing_regions = _find_placeholders_needing_backing(src_slide)
    if not backing_regions:
        return

    src_layout = src_slide.slide_layout
    layout_spTree = src_layout.element.find('.//' + qn('p:cSld')).find(qn('p:spTree'))
    dst_spTree = dst_slide.shapes._spTree

    # Compute scale ratios for cross-dimension copying
    sx = dst_prs.slide_width / src_prs.slide_width
    sy = dst_prs.slide_height / src_prs.slide_height

    shape_tags = {qn('p:sp'), qn('p:grpSp'), qn('p:pic'),
                  qn('p:graphicFrame'), qn('p:cxnSp')}

    # Find max shape id already in dst to avoid conflicts
    max_id = 0
    for cNvPr in dst_spTree.iter(qn('p:cNvPr')):
        try:
            max_id = max(max_id, int(cNvPr.get('id', 0)))
        except ValueError:
            pass

    # Build rId map for layout relationships (images referenced by decorations)
    layout_rid_map = {}
    for rel in src_layout.part.rels.values():
        if 'slideMaster' in rel.reltype or rel.is_external:
            continue
        if 'image' in rel.reltype:
            src_part = rel.target_part
            new_partname = _next_media_partname(src_part.content_type)
            new_part = Part(
                new_partname, src_part.content_type,
                dst_slide.part.package, src_part.blob
            )
            new_rid = dst_slide.part.relate_to(new_part, rel.reltype)
            layout_rid_map[rel.rId] = new_rid
        elif not rel.is_external:
            new_rid = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            layout_rid_map[rel.rId] = new_rid

    # Find insertion point: insert before existing content shapes
    # so decorations render behind content
    first_shape_idx = None
    for i, child in enumerate(dst_spTree):
        if child.tag in shape_tags:
            first_shape_idx = i
            break

    insert_idx = first_shape_idx if first_shape_idx is not None else len(dst_spTree)

    for child in layout_spTree:
        if child.tag not in shape_tags:
            continue

        # Skip placeholder shapes — only copy decorations
        if child.tag == qn('p:sp'):
            ph_info = _get_placeholder_info(child)
            if ph_info is not None:
                continue

        # Only copy if this decoration overlaps with a placeholder that
        # needs a dark backing shape
        deco_rect = _get_shape_rect(child)
        if deco_rect is None:
            continue
        if not any(_rects_overlap(deco_rect, region) for region in backing_regions):
            continue

        new_shape = deepcopy(child)

        # Scale position/size to match destination slide dimensions
        if sx != 1.0 or sy != 1.0:
            _scale_xfrm(new_shape, sx, sy)

        # Assign unique shape ids
        for cNvPr in new_shape.iter(qn('p:cNvPr')):
            max_id += 1
            cNvPr.set('id', str(max_id))

        # Update relationship references (e.g. image embeds in p:pic)
        if layout_rid_map:
            _update_rids_in_tree(new_shape, layout_rid_map)

        dst_spTree.insert(insert_idx, new_shape)
        insert_idx += 1


def _add_empty_placeholders_from_layout(dst_slide, dst_layout):
    """
    Add empty placeholder shapes to the slide for each placeholder in the layout.
    This overrides the layout's placeholder text (which would otherwise show through)
    while keeping non-placeholder shapes (decorations, images) visible.
    """
    dst_spTree = dst_slide.shapes._spTree

    # Find the max shape id already in spTree to avoid conflicts
    max_id = 0
    for cNvPr in dst_spTree.iter(qn('p:cNvPr')):
        try:
            max_id = max(max_id, int(cNvPr.get('id', 0)))
        except ValueError:
            pass

    layout_spTree = dst_layout.element.find('.//' + qn('p:cSld')).find(qn('p:spTree'))
    for sp in layout_spTree.findall(qn('p:sp')):
        ph_info = _get_placeholder_info(sp)
        if ph_info is None:
            continue

        # Create an empty placeholder shape that matches the layout's placeholder
        max_id += 1
        ph_type, ph_idx = ph_info

        # Build minimal sp element with placeholder reference
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                 'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                 'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

        sp_elem = etree.SubElement(dst_spTree, qn('p:sp'))

        # nvSpPr
        nvSpPr = etree.SubElement(sp_elem, qn('p:nvSpPr'))
        cNvPr = etree.SubElement(nvSpPr, qn('p:cNvPr'))
        cNvPr.set('id', str(max_id))
        cNvPr.set('name', f'Empty Placeholder {max_id}')
        cNvSpPr = etree.SubElement(nvSpPr, qn('p:cNvSpPr'))
        sp_locks = etree.SubElement(cNvSpPr, qn('a:spLocks'))
        sp_locks.set('noGrp', '1')
        nvPr = etree.SubElement(nvSpPr, qn('p:nvPr'))
        ph = etree.SubElement(nvPr, qn('p:ph'))
        if ph_type:
            ph.set('type', ph_type)
        if ph_idx:
            ph.set('idx', ph_idx)

        # spPr (empty - inherits position/size from layout)
        etree.SubElement(sp_elem, qn('p:spPr'))

        # txBody with single empty paragraph (overrides layout text)
        txBody = etree.SubElement(sp_elem, qn('p:txBody'))
        bodyPr = etree.SubElement(txBody, qn('a:bodyPr'))
        lstStyle = etree.SubElement(txBody, qn('a:lstStyle'))
        p = etree.SubElement(txBody, qn('a:p'))
        endParaRPr = etree.SubElement(p, qn('a:endParaRPr'))


def _map_placeholders(src_slide, dst_slide, rid_map):
    """
    Map content from source placeholders into template placeholders.
    Non-placeholder shapes and unmatched placeholders are added as extra shapes.
    """
    shape_tags = {qn('p:sp'), qn('p:grpSp'), qn('p:pic'),
                  qn('p:graphicFrame'), qn('p:cxnSp')}

    # Classify source shapes
    src_ph_map = {}     # type_key -> sp_elem
    src_other = []      # non-placeholder shapes

    for child in src_slide.shapes._spTree:
        if child.tag not in shape_tags:
            continue
        ph_info = _get_placeholder_info(child)
        if ph_info:
            type_key = ph_info[0] or f'_idx_{ph_info[1]}'
            if type_key not in src_ph_map:
                src_ph_map[type_key] = child
        else:
            src_other.append(child)

    # Map content to matching template placeholders
    dst_spTree = dst_slide.shapes._spTree
    matched = set()

    for dst_child in list(dst_spTree):
        dst_ph_info = _get_placeholder_info(dst_child)
        if dst_ph_info is None:
            continue
        dst_key = dst_ph_info[0] or f'_idx_{dst_ph_info[1]}'
        if dst_key in src_ph_map and dst_key not in matched:
            _copy_placeholder_content(src_ph_map[dst_key], dst_child)
            matched.add(dst_key)

    # Add non-placeholder shapes from source
    for sp in src_other:
        new_sp = deepcopy(sp)
        _update_rids_in_tree(new_sp, rid_map)
        dst_spTree.append(new_sp)

    # Unmatched source placeholders are intentionally dropped.
    # Template defines the structure — if template doesn't have a matching
    # placeholder, that content is not part of the template's design.


def _extract_theme_colors(prs):
    """
    Extract the color scheme from a presentation's theme.
    Returns dict mapping scheme name -> hex value (uppercase, e.g. 'accent1' -> 'FF5030').
    """
    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_elem = etree.fromstring(rel.target_part.blob)
                clrScheme = theme_elem.find('.//' + qn('a:clrScheme'))
                if clrScheme is None:
                    continue
                colors = {}
                for name in ['dk1', 'dk2', 'lt1', 'lt2',
                             'accent1', 'accent2', 'accent3', 'accent4',
                             'accent5', 'accent6', 'hlink', 'folHlink']:
                    elem = clrScheme.find(qn(f'a:{name}'))
                    if elem is None:
                        continue
                    srgb = elem.find(qn('a:srgbClr'))
                    if srgb is not None:
                        colors[name] = srgb.get('val', '').upper()
                    else:
                        sys_clr = elem.find(qn('a:sysClr'))
                        if sys_clr is not None:
                            colors[name] = sys_clr.get('lastClr', '').upper()
                if colors:
                    return colors
    return {}


def _remap_colors_to_theme(spTree, src_theme_colors):
    """
    Replace hardcoded srgbClr values that match a source theme color
    with the corresponding schemeClr reference.

    This way, colors that were theme-derived in the source will adapt
    to the destination template's theme automatically.
    Only exact matches are replaced — other hardcoded colors are left as-is.
    """
    # Build reverse mapping: hex value -> first matching scheme name
    hex_to_scheme = {}
    for name, hex_val in src_theme_colors.items():
        if hex_val and hex_val not in hex_to_scheme:
            hex_to_scheme[hex_val] = name

    if not hex_to_scheme:
        return

    for srgb_elem in list(spTree.iter(qn('a:srgbClr'))):
        hex_val = srgb_elem.get('val', '').upper()
        if hex_val not in hex_to_scheme:
            continue

        scheme_name = hex_to_scheme[hex_val]
        parent = srgb_elem.getparent()

        # Create schemeClr element with the same namespace
        scheme_elem = parent.makeelement(qn('a:schemeClr'), {'val': scheme_name})

        # Preserve child modifiers (alpha, tint, shade, lumMod, lumOff, etc.)
        for child in list(srgb_elem):
            scheme_elem.append(child)

        parent.replace(srgb_elem, scheme_elem)


def _remap_fonts_to_theme(spTree):
    """
    Replace all hardcoded font references with theme font references.

    Title shapes -> major font (+mj-lt/ea/cs)
    All other shapes/tables -> minor font (+mn-lt/ea/cs)

    This ensures the output uses the template's theme fonts and avoids
    missing font issues on machines that don't have the source fonts.
    """
    font_tags = {
        qn('a:latin'): '-lt',
        qn('a:ea'): '-ea',
        qn('a:cs'): '-cs',
    }
    metric_attrs = ['panose', 'pitchFamily', 'charset']

    def _remap_element(font_elem, prefix):
        typeface = font_elem.get('typeface', '')
        if typeface.startswith('+'):
            return  # Already a theme reference
        suffix = font_tags.get(font_elem.tag)
        if suffix is None:
            return
        font_elem.set('typeface', prefix + suffix)
        for attr in metric_attrs:
            if attr in font_elem.attrib:
                del font_elem.attrib[attr]

    # Track processed elements to avoid double-processing
    processed = set()

    # Process p:sp shapes with title detection
    for sp in spTree.iter(qn('p:sp')):
        prefix = '+mj' if _is_title_shape(sp) else '+mn'
        for tag in font_tags:
            for font_elem in sp.iter(tag):
                _remap_element(font_elem, prefix)
                processed.add(id(font_elem))

    # Process remaining font elements (tables in graphicFrame, etc.) as minor
    for tag in font_tags:
        for font_elem in spTree.iter(tag):
            if id(font_elem) not in processed:
                _remap_element(font_elem, '+mn')


def copy_slide(src_prs, src_slide_index, dst_prs, layout_index=0, apply_template_bg=True,
               remap_fonts=True, remap_colors=True, src_theme_colors=None,
               use_placeholders=False, copy_src_layout_decorations=False):
    """
    Copy a single slide from source presentation to destination.

    Args:
        src_prs: Source Presentation object
        src_slide_index: 0-based index of the slide to copy
        dst_prs: Destination Presentation object
        layout_index: Index of the slide layout in destination to apply (0-based)
        apply_template_bg: If True, use template background; if False, copy source background
        use_placeholders: If True, map source content into template placeholders
            instead of replacing the entire shape tree
        copy_src_layout_decorations: If True, copy non-placeholder decoration shapes
            (bars, circles, lines, images) from the source layout into the slide.
            Useful when the source layout has visual elements (e.g. dark header bars)
            that give context to the content (e.g. white text on dark background).

    Returns:
        The newly created slide in the destination presentation
    """
    src_slide = list(src_prs.slides)[src_slide_index]

    # Create new slide with template layout
    dst_layout = dst_prs.slide_layouts[layout_index]
    dst_slide = dst_prs.slides.add_slide(dst_layout)

    # Step 1: Copy relationships (images, charts, etc.) and build rId mapping
    rid_map = _copy_relationships(src_slide, dst_slide)

    # Step 2: Transfer shapes
    if use_placeholders:
        # Map content into template placeholders; add other shapes alongside
        _map_placeholders(src_slide, dst_slide, rid_map)
    else:
        # Replace entire shape tree with source shapes
        new_spTree = deepcopy(src_slide.shapes._spTree)
        _update_rids_in_tree(new_spTree, rid_map)

        # Classify source placeholders dynamically:
        # - Auto-generated (sldNum, dt, ftr, hdr): REMOVE - template provides these
        # - Content with text (title, body, etc.): KEEP as regular shapes
        # - Content but empty: REMOVE - no point keeping empty boxes
        _AUTO_PH_TYPES = {'sldNum', 'dt', 'ftr', 'hdr'}
        for sp in list(new_spTree.findall(qn('p:sp'))):
            ph_info = _get_placeholder_info(sp)
            if ph_info is None:
                continue
            ph_type = ph_info[0]

            if ph_type in _AUTO_PH_TYPES:
                # Auto-generated placeholder - template provides its own
                new_spTree.remove(sp)
            else:
                # Content placeholder - check if it has actual text
                txBody = sp.find(qn('p:txBody'))
                has_text = False
                if txBody is not None:
                    for t_elem in txBody.iter(qn('a:t')):
                        if t_elem.text and t_elem.text.strip():
                            has_text = True
                            break
                if has_text:
                    # Has content: bake position and convert to regular shape
                    _bake_placeholder_styles(sp, src_slide)
                    _remove_placeholder_ref(sp)
                else:
                    # Empty content placeholder: remove to avoid blank boxes
                    new_spTree.remove(sp)

        # Scale all content shapes if source and destination dimensions differ
        sx = dst_prs.slide_width / src_prs.slide_width
        sy = dst_prs.slide_height / src_prs.slide_height
        if sx != 1.0 or sy != 1.0:
            _scale_xfrm(new_spTree, sx, sy)

        dst_spTree = dst_slide.shapes._spTree
        dst_spTree.getparent().replace(dst_spTree, new_spTree)

        # Invalidate the @lazyproperty cache for 'shapes' so subsequent
        # accesses see the new spTree instead of the stale detached one.
        dst_slide.__dict__.pop('shapes', None)

        # Add empty placeholders to override layout's placeholder text
        # (prevents template text from showing through while keeping decorations)
        _add_empty_placeholders_from_layout(dst_slide, dst_layout)

    # Step 2.5: Copy decoration shapes from source layout (bars, images, etc.)
    if copy_src_layout_decorations:
        _copy_layout_decorations(src_slide, dst_slide, src_prs, dst_prs)

    # Step 3: Remap fonts to template theme fonts
    if remap_fonts:
        _remap_fonts_to_theme(dst_slide.shapes._spTree)

    # Step 4: Remap hardcoded colors to theme color references
    if remap_colors and src_theme_colors:
        _remap_colors_to_theme(dst_slide.shapes._spTree, src_theme_colors)

    # Step 5: Handle background
    if not apply_template_bg:
        # Copy source slide's background (if it has one)
        src_cSld = src_slide.element.find(qn('p:cSld'))
        src_bg = src_cSld.find(qn('p:bg')) if src_cSld is not None else None
        if src_bg is not None:
            dst_cSld = dst_slide.element.find(qn('p:cSld'))
            new_bg = deepcopy(src_bg)
            # Insert background before spTree
            dst_spTree_elem = dst_cSld.find(qn('p:spTree'))
            dst_cSld.insert(list(dst_cSld).index(dst_spTree_elem), new_bg)

    return dst_slide


def copy_slides_to_template(template_path, slide_selections, output_path,
                            layout_index=0, apply_template_bg=True,
                            remap_fonts=True, remap_colors=True,
                            use_placeholders=False,
                            copy_src_layout_decorations=False,
                            slide_size=None):
    """
    Main function: create presentation from template and copy selected slides.

    Args:
        template_path: Path to the template PPTX file
        slide_selections: List of tuples (source_path, slide_indices)
            - source_path: Path to a source PPTX file
            - slide_indices: List of 0-based slide indices to copy
        output_path: Path for the output PPTX file
        layout_index: Template layout index to apply (0-based)
        apply_template_bg: If True, template background is applied to all slides
        copy_src_layout_decorations: If True, copy source layout decoration shapes
        slide_size: Output slide dimensions. Use 'a4' for A4 landscape (297x210mm),
            or a (width, height) tuple in EMU. None keeps template dimensions.

    Returns:
        Path to the saved output file
    """
    print(f"Creating presentation from template: {os.path.basename(template_path)}")
    dst_prs = create_from_template(template_path, slide_size=slide_size)
    print(f"  Available layouts: {[l.name for l in dst_prs.slide_layouts]}")
    print(f"  Using layout [{layout_index}]: \"{dst_prs.slide_layouts[layout_index].name}\"")
    print()

    slide_count = 0
    for src_path, indices in slide_selections:
        src_prs = Presentation(src_path)
        total = len(list(src_prs.slides))
        src_theme_colors = _extract_theme_colors(src_prs) if remap_colors else None
        print(f"Source: {os.path.basename(src_path)} ({total} slides)")

        for idx in indices:
            if idx < 0 or idx >= total:
                print(f"  [SKIP] Slide {idx} out of range (0-{total-1})")
                continue

            copy_slide(src_prs, idx, dst_prs, layout_index, apply_template_bg,
                       remap_fonts, remap_colors, src_theme_colors,
                       use_placeholders, copy_src_layout_decorations)
            slide_count += 1
            print(f"  Copied slide {idx} -> destination slide {slide_count}")

    dst_prs.save(output_path)
    print(f"\nSaved {slide_count} slides to: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# Helper functions for building slide index lists
# ---------------------------------------------------------------------------

def first_n(pptx_path, n):
    """Return indices of the first N slides."""
    total = len(list(Presentation(pptx_path).slides))
    return list(range(min(n, total)))


def last_n(pptx_path, n):
    """Return indices of the last N slides."""
    total = len(list(Presentation(pptx_path).slides))
    return list(range(max(0, total - n), total))


def slide_range(start, end):
    """Return indices from start to end (inclusive, 0-based)."""
    return list(range(start, end + 1))


# ---------------------------------------------------------------------------
# Example / Demo
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

    template_path = os.path.join(BASE_DIR, '202711_東京_第72回日本生殖医学会学術講演会・総会（パシフィコ横浜_ノース）_231227★.pptx')
    upload_path = os.path.join(BASE_DIR, '202602_東京_第56回日本心臓血管外科学術総会（幕張）★.pptx')
    output_path = os.path.join(BASE_DIR, 'output.pptx')

    # Task: copy first 5 slides + slide 6 + last 2 slides from file_upload
    slide_selections = [
        (upload_path, first_n(upload_path, 6)),   # Slides 0,1,2,3,4,thực tế đã apply được những gì từ file template nhỉ5
        (upload_path, [23, 24, 25, 27]),                        # Slide 6 (index 5, 0-based)
        #(upload_path, last_n(upload_path, 2)),     # Slides 9,10
    ]

    copy_slides_to_template(
        template_path=template_path,
        slide_selections=slide_selections,
        output_path=output_path,
        layout_index=0,           # "Cover slide layout"
        apply_template_bg=True,   # Apply template background
        copy_src_layout_decorations=True,  # Copy source layout decorations (bars, images)
        slide_size='a4',          # Output slides in A4 landscape (297x210mm)
    )
